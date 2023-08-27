# ppt部分
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
import pandas as pd


def show_page(now_page, all_slides):
    """
    用于显示指定页面：now_page 内的形状
    :param now_page: slide 的索引值
    :param all_slides: 所有的 slides
    显示 slides索引、形状索引、形状名称；如果有内容则显示内容
    """
    shapes = all_slides[now_page].shapes
    for index, shape in enumerate(shapes):
        print(now_page, index, shape.name)
        if shape.has_text_frame:
            print(shape.text)
        if shape.has_table:
            print(shape.table)


def get_column(df, year, city, column, zone=False):
    """
    获得指定年份、城市（或区域）的指定字段值，仅适用于 九城宏观指标汇总表
    :param df: 数据表的 dataframe
    :param year: 数据年份
    :param city: 数据城市
    :param column: 返回的数据字段
    :param zone: 如果需要的是区域级别，则传入所需的区域名称，否则不用填
    :return: 找到值则返回，没有找到则返回 nan
    """
    # 表名（df）、年份、城市、指标名称、区县（不填则为城市级别）
    df = df[(df.parameter == column) & df.city.str.contains(city) & (df.year == year)].copy()
    df = df[df.district == zone].copy() if zone else df[df.district == '城市'].copy()
    res_num = df.value.to_list()
    return round(res_num[0], 1) if len(res_num) else 'nan'


def get_pos_column(df, column, conditions, perc_col=None, res_df=False):
    """
    更灵活的数据提取方式，可以用于数据集的查找或单条数据的查找，返回的数据集会进行去重
    :param df: 数据表的 dataframe
    :param column: 所需返回的 字段 或 字段列表：'col' or ['col1', 'col2']
    :param conditions: 筛选条件，字典形式：{'key': 'val'}
    :param perc_col: 百分数字段，处理为一位百分比格式 0.05913 -> 5.9
    :param res_df: 如果只有一条数据，且需要返回为 dataframe，则需要设置为 True
    :return: 返回 满足条件的所需字段的列表 或者 满足条件的所需数据的 dataframe
             如果只有一条结果或无结果则返回列表，多条则返回所有数据的 dataframe
    """
    for key, val in conditions.items():
        val = val if isinstance(val, list) else [val]
        cases = df[key] == val[0]
        for item in val[1:]:
            cases = cases | (df[key] == item)
        df = df[cases].copy()
    perc_col = [] if perc_col is None else [perc_col] if isinstance(perc_col, str) else perc_col
    for col in perc_col:
        df[col] = df[col].apply(lambda x: round(x*100, 1))
    df = df[column]
    if df.shape[0] > 1 or res_df:
        df.fillna(0, inplace=True)
        df.drop_duplicates(inplace=True)
        return df
    return df.iloc[0, :].to_list()


def get_pos_n_column(df, num, target, first=True):
    """
    返回 前n条 或者 后n条 所有字段的列表
    :param df: 建议使用 get_pos_column 筛选后的数据 dataframe
    :param num: 所需数据的条目数，即 n
    :param target: 用于排序的字段名称
    :param first: 标识 正序还是逆序，True 表示取大值，False 表示取小值
    :return: 返回所有字段的列表 [('丰台区', '西城区', ...), (1, 2, ...)]
    """
    asc = False if first else True
    df = df.sort_values(target, ascending=asc).reset_index(drop=True)
    return [df.iloc[i].to_list() for i in range(num)]


def cal_data(df, year, city, column, zone=False):
    """
    用于计算同比数据，针对 九城宏观指标汇总数据
    :param df: 数据集的 dataframe
    :param year: 所需数据年份，如：2020
    :param city: 需要计算的城市，如：北京市
    :param column: 需要计算同比的字段，单一值，字符串
    :param zone: 如果需要计算区域同比，则传入区域名称，否则不需要传入
    :return: 返回 上一年数据、当年数据、增加|下降、同比绝对值
    """
    # 表名（df）、年份、城市、指标名称、区县（不填则为城市级别）
    df = df[(df.parameter == column) & df.city.str.contains(city)].copy()
    df = df[df.district == zone].copy() if zone else df[df.district == '城市'].copy()
    cur_year = df[df.year == year].value.to_list()
    cur_year = cur_year[0] if len(cur_year) else 'nan'
    las_year = df[df.year == year-1].value.to_list()
    las_year = las_year[0] if len(las_year) else 'nan'
    res_num = round((cur_year-las_year)/las_year*100, 2) if cur_year != 'nan' and las_year != 'nan' else 'nan'
    num_tag = 'nan' if isinstance(res_num, str) else '增加' if res_num > 0 else '下降'
    res_num = res_num if isinstance(res_num, str) else abs(res_num)
    return las_year, cur_year, num_tag, res_num


def txt_format(item, txt, font, size, color=None, bold=False, center=False):
    """
    用于对PPT中的文字格式进行设置，无返回值
    :param item: 需要填充文字的 paragraphs[n]
    :param txt: 需要填充的文字
    :param font: 文字字体，如："微软雅黑"
    :param size: 文字字号，如：16
    :param color: 文字的rgb颜色，接收数组，如：[255, 100, 100]，默认黑色则不用传
    :param bold: 是否加粗，不加粗则不用传参
    :param center: 是否居中，无需设置则不用传参
    """
    # 元素的paragraph[n]、文本、字体、字号、颜色、是否加粗、是否居中
    item.clear()  # 清除原格式
    if center:
        item.alignment = PP_ALIGN.CENTER
    item = item.add_run()
    item.text = txt
    item.font.bold = bold
    item.font.name = font
    item.font.size = Pt(size)
    if color:
        item.font.color.rgb = RGBColor(color[0], color[1], color[2])


def get_n_row(df, year, city, column, n, first=True):
    """
    针对九城宏观指标汇总数据，返回所需字段的 前n名 或 后n名
    :param df: 接收指标数据集
    :param year: 所需年份
    :param city: 所需城市
    :param column: 用于排序的字段
    :param n: 返回数据的条目数据
    :param first: True 则返回前n名，False 则返回后n名
    :return: 返回城市的区域名称、对应值，如果结果数量小于指定的n，则填充为nan
    """
    # 表名（df）、城市、指标名称、最值提取个数，正序还是倒叙（默认正序）
    df = df[(df.parameter == column) & df.city.str.contains(city) & (df.year == year) & (df.district != '城市')].copy()
    df['value'] = df.value.astype('float')
    index = df.value.nlargest(n).index.to_list() if first else df.value.nsmallest(n).index.to_list()
    res = df.loc[index, ['district', 'value']]
    # 返回区域名称列表、对应区域的数值列表
    districts, values = res.district.to_list(), res.value.to_list()
    while len(districts) < n:
        districts.append('nan')
        values.append('nan')
    return districts, values


def replace_chart_data(chart, df, columns, index=False):
    """
    用于替换图表中的数据源，无返回值
    :param chart: 指定图表：shapes[n].chart
    :param df: 用于填充的数据，可使用：get_pos_column 返回的内容
    :param columns: 数据名称及填充顺序
    :param index: 图表的索引字段
    使用示例：
    chart = slides[3].shapes[2].chart                                      # 使用变量保存 chart
    r = get_pos_column(data, ['year', 'parameter', 'value'],               # 筛选数据
                       {'city': '玉林市', 'district': '城市',
                        'parameter': ['人均地区生产总值', '常住人口'],
                        'year': [2016, 2017, 2018, 2019]})
    r = pd.pivot(r, index='year', columns='parameter', values='value')     # 整理数据格式
    replace_chart_data(chart, r, ['常住人口', '人均地区生产总值'])              # 填充数据
    title = chart.chart_title.text_frame.paragraphs[0]                     # 修改图表的标题
    txt_format(title, '北京户籍人口变化和人均GDP', '宋体', 18, bold=True)
    """
    chart_data = ChartData()
    chart_data.categories = df.index.to_list() if not index else df[index].to_list()
    for col in columns:
        chart_data.add_series(col, df[col].to_list())
    chart.replace_data(chart_data)
    chart.value_axis.major_unit = None
    chart.value_axis.minor_unit = None
